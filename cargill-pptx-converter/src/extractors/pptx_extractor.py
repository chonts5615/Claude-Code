"""
PowerPoint (.pptx) content extractor.
"""

from pathlib import Path
from uuid import uuid4

from src.extractors.base import BaseExtractor
from src.schemas.content import (
    ContentBlock,
    ContentType,
    DocumentMetadata,
    ExtractedContent,
)


class PptxExtractor(BaseExtractor):
    """Extract content from PowerPoint presentations."""

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def extract(self, file_path: Path) -> ExtractedContent:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(str(file_path))
        blocks: list[ContentBlock] = []
        warnings: list[str] = []
        word_count = 0
        has_images = False
        has_tables = False
        has_charts = False

        for slide_idx, slide in enumerate(prs.slides):
            slide_num = slide_idx + 1

            for shape in slide.shapes:
                # Text frames (titles, body text, text boxes)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue

                        word_count += len(text.split())
                        block_id = f"pptx_s{slide_num}_{uuid4().hex[:6]}"

                        # Determine if this is a title
                        if shape.shape_id == 0 or (
                            hasattr(shape, "placeholder_format")
                            and shape.placeholder_format is not None
                            and shape.placeholder_format.idx in (0, 1)
                        ):
                            # Title or subtitle placeholder
                            level = 1 if (
                                hasattr(shape, "placeholder_format")
                                and shape.placeholder_format is not None
                                and shape.placeholder_format.idx == 0
                            ) else 2

                            blocks.append(ContentBlock(
                                block_id=block_id,
                                content_type=ContentType.HEADING,
                                text=text,
                                level=level,
                                metadata={"slide": slide_num},
                            ))
                        elif para.level > 0:
                            # Bulleted content
                            if blocks and blocks[-1].content_type == ContentType.BULLET_LIST and \
                               blocks[-1].metadata.get("slide") == slide_num:
                                blocks[-1].items.append(text)
                            else:
                                blocks.append(ContentBlock(
                                    block_id=block_id,
                                    content_type=ContentType.BULLET_LIST,
                                    items=[text],
                                    metadata={"slide": slide_num},
                                ))
                        else:
                            blocks.append(ContentBlock(
                                block_id=block_id,
                                content_type=ContentType.PARAGRAPH,
                                text=text,
                                metadata={"slide": slide_num},
                            ))

                # Tables
                if shape.has_table:
                    has_tables = True
                    table = shape.table
                    rows = []
                    headers = []

                    for row_idx, row in enumerate(table.rows):
                        row_data = [cell.text.strip() for cell in row.cells]
                        if row_idx == 0:
                            headers = row_data
                        else:
                            rows.append(row_data)

                    block_id = f"pptx_tbl_s{slide_num}_{uuid4().hex[:6]}"
                    blocks.append(ContentBlock(
                        block_id=block_id,
                        content_type=ContentType.TABLE,
                        table_data=rows,
                        table_headers=headers,
                        metadata={"slide": slide_num},
                    ))

                # Images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_images = True
                    try:
                        image = shape.image
                        block_id = f"pptx_img_s{slide_num}_{uuid4().hex[:6]}"
                        blocks.append(ContentBlock(
                            block_id=block_id,
                            content_type=ContentType.IMAGE,
                            image_data=image.blob,
                            image_format=image.content_type.split("/")[-1],
                            metadata={"slide": slide_num},
                        ))
                    except Exception as e:
                        warnings.append(f"Slide {slide_num}: Failed to extract image - {e}")

                # Charts
                if shape.has_chart:
                    has_charts = True
                    try:
                        chart_data = self._extract_chart_data(shape.chart)
                        block_id = f"pptx_chart_s{slide_num}_{uuid4().hex[:6]}"
                        blocks.append(ContentBlock(
                            block_id=block_id,
                            content_type=ContentType.CHART_DATA,
                            chart_data=chart_data,
                            metadata={"slide": slide_num},
                        ))
                    except Exception as e:
                        warnings.append(f"Slide {slide_num}: Failed to extract chart - {e}")

        # Infer title
        title = ""
        for block in blocks:
            if block.content_type == ContentType.HEADING and block.level == 1:
                title = block.text or ""
                break

        metadata = DocumentMetadata(
            title=title,
            page_count=len(prs.slides),
            word_count=word_count,
            has_images=has_images,
            has_tables=has_tables,
            has_charts=has_charts,
        )

        return ExtractedContent(
            source_file=str(file_path),
            source_format="pptx",
            blocks=blocks,
            metadata=metadata,
            extraction_warnings=warnings,
        )

    def _extract_chart_data(self, chart) -> dict:
        """Extract data from a chart object."""
        chart_data = {
            "chart_type": str(chart.chart_type) if chart.chart_type else "unknown",
            "categories": [],
            "series": [],
        }

        try:
            plot = chart.plots[0]
            # Get categories
            if hasattr(plot, "categories") and plot.categories:
                chart_data["categories"] = list(plot.categories)

            # Get series
            for series in plot.series:
                series_data = {
                    "name": str(series.tx) if hasattr(series, "tx") else "",
                    "values": list(series.values) if hasattr(series, "values") else [],
                }
                chart_data["series"].append(series_data)
        except Exception:
            pass

        return chart_data
