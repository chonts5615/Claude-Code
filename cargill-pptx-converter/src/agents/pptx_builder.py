"""
Agent 7: PPTX Builder Agent

Renders the final PowerPoint presentation file using python-pptx.
This is the core rendering engine that translates SlideSpec objects into
actual PowerPoint slides with full Cargill brand compliance.
"""

import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.agents.base import BaseAgent
from src.brand.constants import (
    BODY_BASE,
    BODY_SM,
    CARGILL_LEAF_GREEN,
    CORNER_RADIUS_MD,
    DEEP_GREEN,
    FOOTER_COLOR,
    FOOTER_FONT_SIZE,
    FOOTER_TEXT,
    HEADING_FONT,
    HEADING_FONT_FALLBACK,
    BODY_FONT,
    BODY_FONT_FALLBACK,
    HEADING_SM,
    NEUTRAL_200,
    NEUTRAL_300,
    NEUTRAL_700,
    NEUTRAL_1000,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    SOFT_GREEN,
    WHITE,
    WHITE_GREEN,
)
from src.brand.palette import hex_to_rgb
from src.schemas.slide import SlideLayout, SlideSpec
from src.schemas.run_state import RunState


class PptxBuilderAgent(BaseAgent):
    """Render the final PowerPoint presentation."""

    def __init__(self):
        super().__init__("S7", "PPTX Rendering")

    def execute(self, state: RunState) -> RunState:
        state.current_step = self.agent_id
        self.logger.info("Building PowerPoint presentation")

        if not state.presentation_plan:
            self._add_flag(state, "CRITICAL", "No presentation plan for PPTX building")
            return state

        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        use_brand_fonts = state.config.use_brand_fonts

        for slide_spec in state.presentation_plan.slides:
            try:
                self._render_slide(prs, slide_spec, state.config, use_brand_fonts)
            except Exception as e:
                self._add_flag(
                    state, "WARNING",
                    f"Failed to render slide {slide_spec.slide_number}: {e}"
                )

        # Save output
        output_path = self._resolve_output_path(state)
        prs.save(str(output_path))
        state.output_file = str(output_path)
        state.artifacts.output_pptx_path = str(output_path)

        self.logger.info(f"Saved presentation: {output_path}")
        return state

    def _resolve_output_path(self, state: RunState) -> Path:
        """Determine output file path."""
        if state.inputs.output_path:
            return Path(state.inputs.output_path)

        input_path = Path(state.inputs.input_file)
        output_dir = input_path.parent
        output_name = f"{input_path.stem}_Cargill_Branded.pptx"
        return output_dir / output_name

    def _render_slide(self, prs, spec: SlideSpec, config, use_brand_fonts: bool):
        """Render a single slide based on its specification."""
        # Use blank layout
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Apply background
        self._apply_background(slide, spec)

        # Render based on layout type
        layout_renderers = {
            SlideLayout.TITLE_HERO: self._render_hero_slide,
            SlideLayout.BASIC_HERO: self._render_hero_slide,
            SlideLayout.SECTION_HEADER: self._render_section_header,
            SlideLayout.CONTENT: self._render_content_slide,
            SlideLayout.CONTENT_LEFT_HEADLINE: self._render_content_slide,
            SlideLayout.CONTENT_CENTERED: self._render_centered_slide,
            SlideLayout.BULLET_LIST: self._render_bullet_slide,
            SlideLayout.TABLE: self._render_table_slide,
            SlideLayout.CHART: self._render_chart_slide,
            SlideLayout.IMAGE_WITH_TEXT: self._render_image_slide,
            SlideLayout.IMAGE_FULL: self._render_image_slide,
            SlideLayout.SIMPLE_STATISTICS: self._render_stats_slide,
            SlideLayout.STATISTIC_CARDS: self._render_stat_cards_slide,
            SlideLayout.STATS_WITH_HEADLINE: self._render_stats_slide,
            SlideLayout.TWO_COLUMN: self._render_two_column_slide,
            SlideLayout.THREE_COLUMN_CARDS: self._render_three_column_slide,
            SlideLayout.CLOSING: self._render_closing_slide,
            SlideLayout.BLANK: lambda s, sp, f: None,
        }

        renderer = layout_renderers.get(spec.layout, self._render_content_slide)
        renderer(slide, spec, use_brand_fonts)

        # Add footer (except on hero and closing slides)
        if config.include_footer and spec.template_category not in ("hero", "closing"):
            self._add_footer(slide, spec, use_brand_fonts)

        # Add logo (except on closing slides with green background)
        if config.include_logo and spec.layout != SlideLayout.CLOSING:
            self._add_logo_placeholder(slide, spec)

    def _apply_background(self, slide, spec: SlideSpec):
        """Apply background color or fill to slide."""
        bg_color = spec.color_scheme.background
        if bg_color and bg_color != "#FFFFFF":
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(bg_color)

    def _get_heading_font(self, use_brand_fonts: bool) -> str:
        return HEADING_FONT if use_brand_fonts else HEADING_FONT_FALLBACK

    def _get_body_font(self, use_brand_fonts: bool) -> str:
        return BODY_FONT if use_brand_fonts else BODY_FONT_FALLBACK

    # =========================================================================
    # SLIDE RENDERERS
    # =========================================================================

    def _render_hero_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a hero/title slide."""
        padding = spec.spacing.content_padding

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), Inches(2.0),
                Inches(10.0), Inches(2.5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.alignment = PP_ALIGN.LEFT

        # Subtitle/body
        body_text = self._get_element_content(spec, "body")
        if body_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), Inches(4.8),
                Inches(8.0), Inches(1.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_text
            p.font.size = Pt(spec.typography.body_size)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.body_text)
            p.alignment = PP_ALIGN.LEFT

        # Decorative accent bar at bottom
        if spec.graphical_device == "leaf_with_stripe":
            accent_color = spec.color_scheme.accent or "#57D1FF"
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0), Inches(7.0),
                SLIDE_WIDTH, Inches(0.5),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(accent_color)
            shape.line.fill.background()

    def _render_section_header(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a section header/divider slide."""
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5),
                Inches(10.0), Inches(2.5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.auto_size = None
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.alignment = PP_ALIGN.CENTER

        # Green accent line
        shape = slide.shapes.add_shape(
            1, Inches(5.5), Inches(5.3), Inches(2.333), Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARGILL_LEAF_GREEN
        shape.line.fill.background()

    def _render_content_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a standard content slide with title and body."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.font.bold = False
            y_pos = Inches(1.8)

        # Body
        body_text = self._get_element_content(spec, "body")
        if body_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(4.5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_text
            p.font.size = Pt(spec.typography.body_size)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.body_text)
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.LEFT

    def _render_centered_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a centered content slide."""
        body_text = self._get_element_content(spec, "body")
        if body_text:
            txBox = slide.shapes.add_textbox(
                Inches(2.0), Inches(2.0),
                Inches(9.333), Inches(3.5),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_text
            p.font.size = Pt(24)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER

    def _render_bullet_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a bullet list slide."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            y_pos = Inches(1.8)

        # Bullets
        bullet_element = self._get_element(spec, "bullet")
        if bullet_element and bullet_element.items:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(5.0),
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, item in enumerate(bullet_element.items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(spec.typography.body_size)
                p.font.name = self._get_body_font(use_brand_fonts)
                p.font.color.rgb = hex_to_rgb(spec.color_scheme.body_text)
                p.space_after = Pt(8)
                p.level = 0
                # Add bullet character
                pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", {}
                )
                buChar.set("char", "\u2022")
                pPr.append(buChar)

    def _render_table_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a table slide."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            y_pos = Inches(1.8)

        # Find table element
        table_element = self._get_element(spec, "table")
        if not table_element or not table_element.table_data:
            return

        headers = table_element.table_headers or []
        data = table_element.table_data

        # Calculate dimensions
        num_rows = len(data) + (1 if headers else 0)
        num_cols = len(headers) if headers else (len(data[0]) if data else 1)
        num_rows = max(num_rows, 1)
        num_cols = max(num_cols, 1)

        table_width = Inches(min(11.0, num_cols * 2.0))
        row_height = Inches(0.4)
        table_height = row_height * num_rows

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(padding), y_pos,
            table_width, min(table_height, Inches(5.0)),
        )
        table = table_shape.table

        body_font = self._get_body_font(use_brand_fonts)

        # Style header row
        row_idx = 0
        if headers:
            for col_idx, header in enumerate(headers[:num_cols]):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARGILL_LEAF_GREEN

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.name = body_font
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
            row_idx = 1

        # Style data rows
        for i, row_data in enumerate(data):
            table_row = i + row_idx
            if table_row >= num_rows:
                break

            bg = WHITE_GREEN if i % 2 == 0 else WHITE

            for col_idx, cell_value in enumerate(row_data[:num_cols]):
                cell = table.cell(table_row, col_idx)
                cell.text = str(cell_value) if cell_value else ""
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = body_font
                    paragraph.font.color.rgb = NEUTRAL_1000
                    paragraph.alignment = PP_ALIGN.LEFT

        # Style table borders
        for row in range(num_rows):
            for col in range(num_cols):
                cell = table.cell(row, col)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)

    def _render_chart_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a chart slide (embeds pre-rendered chart image)."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            y_pos = Inches(1.8)

        # Find chart image
        chart_element = self._get_element(spec, "chart")
        if chart_element and chart_element.image_data:
            image_stream = io.BytesIO(chart_element.image_data)
            slide.shapes.add_picture(
                image_stream,
                Inches(padding + 0.5), y_pos,
                Inches(10.0), Inches(5.0),
            )
        elif chart_element:
            # Fallback: add placeholder text
            txBox = slide.shapes.add_textbox(
                Inches(2.0), Inches(3.0),
                Inches(9.0), Inches(1.0),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = "[Chart data available - rendering pending]"
            p.font.size = Pt(14)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = NEUTRAL_700
            p.alignment = PP_ALIGN.CENTER

    def _render_image_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a slide with an image."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            y_pos = Inches(1.8)

        # Image
        image_element = self._get_element(spec, "image")
        if image_element and image_element.image_data:
            image_stream = io.BytesIO(image_element.image_data)
            try:
                slide.shapes.add_picture(
                    image_stream,
                    Inches(padding), y_pos,
                    Inches(10.0), Inches(5.0),
                )
            except Exception:
                pass

    def _render_stats_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a statistics slide with numbers in a row."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.alignment = PP_ALIGN.CENTER
            y_pos = Inches(2.5)

        # Stats
        stat_elements = [e for e in spec.elements if e.element_type == "stat"]
        if not stat_elements:
            return

        num_stats = len(stat_elements)
        stat_width = 10.0 / max(num_stats, 1)
        start_x = (13.333 - stat_width * num_stats) / 2

        for i, stat in enumerate(stat_elements):
            x = Inches(start_x + i * stat_width)

            # Stat value
            txBox = slide.shapes.add_textbox(
                x, y_pos, Inches(stat_width), Inches(1.2),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stat.stat_value or ""
            p.font.size = Pt(spec.typography.stat_size)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.primary)
            p.alignment = PP_ALIGN.CENTER

            # Stat label
            txBox = slide.shapes.add_textbox(
                x, Inches(y_pos.inches + 1.4), Inches(stat_width), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stat.stat_label or ""
            p.font.size = Pt(14)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.body_text)
            p.alignment = PP_ALIGN.CENTER

    def _render_stat_cards_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render statistics in card layout."""
        padding = spec.spacing.content_padding
        y_pos = Inches(0.8)

        # Title
        title_text = self._get_element_content(spec, "title")
        if title_text:
            txBox = slide.shapes.add_textbox(
                Inches(padding), y_pos,
                Inches(11.0), Inches(0.8),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = Pt(spec.typography.headline_size)
            p.font.name = self._get_heading_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
            p.alignment = PP_ALIGN.CENTER
            y_pos = Inches(2.0)

        # Stat cards in 2x2 grid
        stat_elements = [e for e in spec.elements if e.element_type == "stat"]
        if not stat_elements:
            return

        card_w = Inches(4.5)
        card_h = Inches(2.2)
        gap = Inches(0.5)
        grid_w = card_w * 2 + gap
        start_x = (SLIDE_WIDTH - grid_w) / 2

        for i, stat in enumerate(stat_elements[:4]):
            col = i % 2
            row = i // 2
            x = start_x + col * (card_w + gap)
            y = y_pos + row * (card_h + gap)

            # Card background
            card = slide.shapes.add_shape(
                1,  # Rectangle
                x, y, card_w, card_h,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.fill.background()
            card.shadow.inherit = False

            # Stat value
            txBox = slide.shapes.add_textbox(
                x + Inches(0.3), y + Inches(0.4),
                card_w - Inches(0.6), Inches(1.0),
            )
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = stat.stat_value or ""
            p.font.size = Pt(36)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.bold = True
            p.font.color.rgb = CARGILL_LEAF_GREEN
            p.alignment = PP_ALIGN.LEFT

            # Stat label
            txBox = slide.shapes.add_textbox(
                x + Inches(0.3), y + Inches(1.4),
                card_w - Inches(0.6), Inches(0.6),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = stat.stat_label or ""
            p.font.size = Pt(14)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = NEUTRAL_700
            p.alignment = PP_ALIGN.LEFT

    def _render_two_column_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a two-column layout slide."""
        self._render_content_slide(slide, spec, use_brand_fonts)

    def _render_three_column_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render a three-column card layout slide."""
        self._render_content_slide(slide, spec, use_brand_fonts)

    def _render_closing_slide(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Render the closing/thank you slide."""
        # Title centered
        title_text = self._get_element_content(spec, "title") or "Thank You"
        txBox = slide.shapes.add_textbox(
            Inches(2.0), Inches(2.0),
            Inches(9.333), Inches(2.0),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(spec.typography.headline_size)
        p.font.name = self._get_heading_font(use_brand_fonts)
        p.font.color.rgb = hex_to_rgb(spec.color_scheme.headline)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        body_text = self._get_element_content(spec, "body")
        if body_text:
            txBox = slide.shapes.add_textbox(
                Inches(2.0), Inches(4.2),
                Inches(9.333), Inches(1.0),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = body_text
            p.font.size = Pt(20)
            p.font.name = self._get_body_font(use_brand_fonts)
            p.font.color.rgb = hex_to_rgb(spec.color_scheme.body_text)
            p.alignment = PP_ALIGN.CENTER

        # Accent line
        line_width = Inches(3.0)
        line_x = (SLIDE_WIDTH - line_width) / 2
        shape = slide.shapes.add_shape(
            1, line_x, Inches(5.5), line_width, Inches(0.04),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.fill.background()

    # =========================================================================
    # HELPER METHODS
    # =========================================================================

    def _get_element_content(self, spec: SlideSpec, element_type: str) -> str:
        """Get content text from a specific element type."""
        for element in spec.elements:
            if element.element_type == element_type and element.content:
                return element.content
        return ""

    def _get_element(self, spec: SlideSpec, element_type: str):
        """Get a specific element from slide spec."""
        for element in spec.elements:
            if element.element_type == element_type:
                return element
        return None

    def _add_footer(self, slide, spec: SlideSpec, use_brand_fonts: bool):
        """Add footer text to slide."""
        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.0),
            Inches(5.0), Inches(0.3),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = FOOTER_TEXT
        p.font.size = FOOTER_FONT_SIZE
        p.font.name = self._get_body_font(use_brand_fonts)
        p.font.color.rgb = NEUTRAL_700
        p.alignment = PP_ALIGN.LEFT

    def _add_logo_placeholder(self, slide, spec: SlideSpec):
        """Add a logo text placeholder (since actual logo file may not be available)."""
        # Use text as placeholder since logo image may not exist
        txBox = slide.shapes.add_textbox(
            Inches(11.0), Inches(0.2),
            Inches(2.0), Inches(0.5),
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "CARGILL"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.RIGHT

        # Set color based on background
        if spec.template_category in ("hero", "closing"):
            p.font.color.rgb = WHITE
        else:
            p.font.color.rgb = CARGILL_LEAF_GREEN
